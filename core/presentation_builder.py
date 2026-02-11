"""
テンプレート設定対応版 PresentationBuilder

任意のテンプレートに対応し、設定ファイルで柔軟に管理
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from typing import Union, List, Dict, Any
import json


class TemplateConfig:
    """テンプレート設定クラス"""
    
    def __init__(self, config_dict: Dict[str, Any] = None):
        """
        Args:
            config_dict: テンプレート設定辞書
        """
        self.config = config_dict or self._default_config()
    
    @staticmethod
    def _default_config():
        """デフォルト設定（一般的なPowerPointテンプレート）"""
        return {
            'title_slide': 0,
            'content_slide': 1,
            'two_column': 3,
            'blank': 6
        }
    
    @classmethod
    def from_file(cls, config_path: str):
        """JSONファイルから設定を読み込む"""
        with open(config_path, 'r', encoding='utf-8') as f:
            config_dict = json.load(f)
        return cls(config_dict)
    
    def save(self, config_path: str):
        """設定をJSONファイルに保存"""
        with open(config_path, 'w', encoding='utf-8') as f:
            json.dump(self.config, f, ensure_ascii=False, indent=2)
    
    def get_layout_index(self, layout_type: str, default: int = 0) -> int:
        """レイアウトタイプからインデックスを取得"""
        return self.config.get(layout_type, default)


class PresentationBuilder:
    """テンプレートベースのプレゼンテーション作成クラス（設定対応版）"""
    
    def __init__(self, template_path: str, config: TemplateConfig = None):
        """
        Args:
            template_path: テンプレートファイルのパス
            config: テンプレート設定（省略時はデフォルト）
        """
        self.template_path = template_path
        self.prs = Presentation(template_path)
        self.config = config or TemplateConfig()
        self._clear_slides()
        self._build_layout_index()
    
    def _clear_slides(self):
        """既存スライドをすべて削除"""
        while len(self.prs.slides) > 0:
            rId = self.prs.slides._sldIdLst[0].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[0]
    
    def _build_layout_index(self):
        """レイアウトのインデックスを作成"""
        self.layouts = {}
        for idx, layout in enumerate(self.prs.slide_layouts):
            self.layouts[layout.name] = idx
            self.layouts[idx] = idx
    
    def list_layouts(self) -> List[Dict[str, Any]]:
        """利用可能なレイアウト一覧を取得"""
        layouts_info = []
        for idx, layout in enumerate(self.prs.slide_layouts):
            placeholders = []
            for shape in layout.placeholders:
                placeholders.append({
                    'idx': shape.placeholder_format.idx,
                    'name': shape.name,
                    'type': str(shape.placeholder_format.type)
                })
            
            layouts_info.append({
                'index': idx,
                'name': layout.name,
                'placeholders': placeholders
            })
        
        return layouts_info
    
    def print_layouts(self):
        """レイアウト一覧を標準出力に表示"""
        print("=== 利用可能なレイアウト ===")
        for layout_info in self.list_layouts():
            print(f"[{layout_info['index']}] {layout_info['name']}")
            for ph in layout_info['placeholders']:
                print(f"    - {ph['name']} (idx={ph['idx']}, type={ph['type']})")
        print()
    
    def generate_config_template(self, output_path: str = None):
        """
        現在のテンプレートから設定ファイルのテンプレートを生成
        
        Args:
            output_path: 出力先パス（省略時は標準出力）
        """
        config = {
            'title_slide': 0,
            'content_slide': 1,
            'two_column': None,
            'blank': None,
            'note': 'レイアウトインデックスを確認して適切な値を設定してください'
        }
        
        if output_path:
            with open(output_path, 'w', encoding='utf-8') as f:
                json.dump(config, f, ensure_ascii=False, indent=2)
            print(f"設定テンプレートを保存: {output_path}")
        else:
            print(json.dumps(config, ensure_ascii=False, indent=2))
    
    def add_slide(self, layout_name_or_index: Union[str, int] = 0):
        """
        新しいスライドを追加
        
        Args:
            layout_name_or_index: レイアウト名、インデックス、または設定キー
        """
        if isinstance(layout_name_or_index, str):
            # 設定キーまたはレイアウト名として扱う
            layout_idx = self.config.get_layout_index(
                layout_name_or_index, 
                self.layouts.get(layout_name_or_index, 0)
            )
        else:
            layout_idx = layout_name_or_index
        
        slide_layout = self.prs.slide_layouts[layout_idx]
        return self.prs.slides.add_slide(slide_layout)
    
    def add_title_slide(self, title: str, subtitle: str = None):
        """タイトルスライドを追加"""
        slide = self.add_slide('title_slide')
        slide.shapes.title.text = title
        
        if subtitle:
            for shape in slide.placeholders:
                if shape.placeholder_format.idx != 0:
                    shape.text = subtitle
                    break
        
        return slide
    
    def add_content_slide(self, title: str, content: List[str] = None):
        """コンテンツスライドを追加"""
        slide = self.add_slide('content_slide')
        slide.shapes.title.text = title
        
        if content:
            for shape in slide.placeholders:
                if shape.placeholder_format.idx != 0 and hasattr(shape, 'text_frame'):
                    text_frame = shape.text_frame
                    text_frame.clear()
                    
                    for i, item in enumerate(content):
                        if i == 0:
                            text_frame.text = item
                        else:
                            p = text_frame.add_paragraph()
                            p.text = item
                            p.level = 0
                    break
        
        return slide
    
    def save(self, output_path: str):
        """プレゼンテーションを保存"""
        self.prs.save(output_path)
    
    @property
    def slide_count(self) -> int:
        """現在のスライド数を取得"""
        return len(self.prs.slides)


def create_simple_presentation(template_path: str, output_path: str, 
                              slides_data: List[Dict[str, Any]],
                              config_path: str = None) -> PresentationBuilder:
    """
    Helper function to create a presentation from structured data
    
    Args:
        template_path: Path to template file
        output_path: Path for output file
        slides_data: List of slide specifications
        config_path: Optional path to config JSON
    
    Returns:
        PresentationBuilder instance
    """
    # Load config if provided
    config = TemplateConfig.from_file(config_path) if config_path else None
    
    # Initialize builder
    builder = PresentationBuilder(template_path, config)
    
    # Create slides
    for slide_spec in slides_data:
        slide_type = slide_spec.get('type', 'content')
        
        if slide_type == 'title':
            builder.add_title_slide(
                slide_spec['title'],
                slide_spec.get('subtitle')
            )
        elif slide_type == 'content':
            builder.add_content_slide(
                slide_spec['title'],
                slide_spec.get('content', [])
            )
        elif slide_type == 'two_column':
            builder.add_two_column_slide(
                slide_spec['title'],
                slide_spec.get('left', []),
                slide_spec.get('right', [])
            )
        elif slide_type == 'custom':
            builder.add_slide(slide_spec.get('layout', 0))
        else:
            # Default to content slide
            builder.add_content_slide(
                slide_spec.get('title', 'Untitled'),
                slide_spec.get('content', [])
            )
    
    # Save
    builder.save(output_path)
    
    return builder


# 使用例
def example_with_config():
    """設定ファイルを使った例"""
    
    # 1. 新しいテンプレート用の設定を生成
    builder = PresentationBuilder('template.pptx')
    builder.print_layouts()
    builder.generate_config_template('template_config.json')
    
    # 2. 設定を編集（手動またはプログラム）
    config = TemplateConfig({
        'title_slide': 0,
        'content_slide': 4,  # 実際のレイアウトに合わせて調整
        'two_column': 7,
        'blank': 13
    })
    
    # 3. 設定を使ってプレゼンテーション作成
    builder = PresentationBuilder('template.pptx', config)
    builder.add_title_slide("タイトル")
    builder.add_content_slide("内容", ["項目1", "項目2"])
    builder.save('output.pptx')


def example_simple():
    """シンプルな使い方（設定なし）"""
    
    # デフォルト設定で動作
    builder = PresentationBuilder('template.pptx')
    
    # レイアウトを確認
    builder.print_layouts()
    
    # インデックス直接指定でも使える
    slide = builder.add_slide(0)
    slide.shapes.title.text = "タイトル"
    
    builder.save('output.pptx')


if __name__ == '__main__':
    example_simple()
