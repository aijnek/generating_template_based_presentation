#!/usr/bin/env python3
"""
Template-Conformant PPTX Helper Functions
==========================================
テンプレート準拠スライド作成のためのヘルパー関数集
"""

from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE


def duplicate_slide(prs, slide_idx):
    """
    指定されたスライドを複製
    
    Args:
        prs: Presentationオブジェクト
        slide_idx: 複製元のスライドインデックス（0始まり）
    
    Returns:
        新しく作成されたスライドオブジェクト
    """
    source = prs.slides[slide_idx]
    
    # 同じレイアウトで新規スライドを作成
    blank_slide = prs.slides.add_slide(source.slide_layout)
    
    # すべての要素をコピー
    for shape in source.shapes:
        el = shape.element
        newel = deepcopy(el)
        blank_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
    
    return blank_slide


def is_title_shape(shape):
    """
    シェイプがタイトルかどうかを判定
    
    Args:
        shape: 判定対象のシェイプ
    
    Returns:
        bool: タイトルの場合True
    """
    # プレースホルダーの場合、タイプで判定
    if shape.is_placeholder:
        try:
            ph_type = shape.placeholder_format.type
            # TITLE (1), SUBTITLE (3), CENTER_TITLE (14)
            return ph_type in [1, 3, 14]
        except:
            pass
    
    # プレースホルダーでない場合、位置とサイズで判定
    # 上部1/3にある、ある程度大きなテキストボックス
    if shape.has_text_frame:
        slide_height = shape.part.slide.presentation.slide_height
        if shape.top < slide_height / 3 and shape.height > Inches(0.5):
            return True
    
    return False


def is_content_shape(shape):
    """
    シェイプがメインコンテンツ（本文）かどうかを判定
    
    Args:
        shape: 判定対象のシェイプ
    
    Returns:
        bool: メインコンテンツの場合True
    """
    # プレースホルダーの場合、タイプで判定
    if shape.is_placeholder:
        try:
            ph_type = shape.placeholder_format.type
            # BODY (2), OBJECT (7)
            return ph_type in [2, 7]
        except:
            pass
    
    # タイトルでなく、テキストを持つシェイプ
    if shape.has_text_frame and not is_title_shape(shape):
        return True
    
    return False


def get_text_shapes(slide):
    """
    スライド内のテキストシェイプを取得
    
    Args:
        slide: 対象スライド
    
    Returns:
        dict: {'title': shape, 'content': [shapes], 'other': [shapes]}
    """
    result = {
        'title': None,
        'content': [],
        'other': []
    }
    
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        if is_title_shape(shape):
            result['title'] = shape
        elif is_content_shape(shape):
            result['content'].append(shape)
        else:
            result['other'].append(shape)
    
    return result


def replace_slide_text(slide, title=None, content=None):
    """
    スライドのテキストを置き換え
    
    Args:
        slide: 対象スライド
        title: 新しいタイトルテキスト（Noneの場合は変更なし）
        content: 新しいコンテンツテキスト（Noneの場合は変更なし）
                 文字列またはリスト（箇条書きの場合）
    
    Returns:
        bool: 成功した場合True
    """
    shapes = get_text_shapes(slide)
    
    # タイトルの置き換え
    if title is not None and shapes['title'] is not None:
        shapes['title'].text = title
    
    # コンテンツの置き換え
    if content is not None and shapes['content']:
        content_shape = shapes['content'][0]  # 最初のコンテンツシェイプ
        
        # リストの場合は箇条書きとして設定
        if isinstance(content, list):
            text_frame = content_shape.text_frame
            text_frame.clear()
            
            for i, item in enumerate(content):
                p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                p.text = item
                p.level = 0  # インデントレベル
        else:
            # 文字列の場合はそのまま設定
            content_shape.text = content
    
    return True


def analyze_layout_usage(prs):
    """
    レイアウトの使用状況を分析
    
    Args:
        prs: Presentationオブジェクト
    
    Returns:
        tuple: (layout_usage dict, issues list)
    """
    layout_usage = {}
    consecutive_runs = []
    current_run = {'layout': None, 'count': 0, 'start_idx': 0}
    
    for i, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name
        
        # カウント
        layout_usage[layout_name] = layout_usage.get(layout_name, 0) + 1
        
        # 連続使用のチェック
        if layout_name == current_run['layout']:
            current_run['count'] += 1
        else:
            if current_run['count'] > 0:
                consecutive_runs.append(current_run.copy())
            current_run = {
                'layout': layout_name,
                'count': 1,
                'start_idx': i
            }
    
    # 最後のrunを追加
    if current_run['count'] > 0:
        consecutive_runs.append(current_run)
    
    # 問題のある連続使用を報告
    issues = []
    for run in consecutive_runs:
        if run['count'] >= 3:
            issues.append({
                'layout': run['layout'],
                'count': run['count'],
                'start_slide': run['start_idx'] + 1,  # 1始まりに変換
                'message': f"スライド{run['start_idx']+1}から: '{run['layout']}'が{run['count']}回連続使用"
            })
    
    return layout_usage, issues


def check_text_overflow(slide):
    """
    スライド内のテキストがはみ出していないかチェック
    
    Args:
        slide: 対象スライド
    
    Returns:
        list: 問題のあるシェイプのリスト
    """
    issues = []
    
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        text_frame = shape.text_frame
        
        # テキストがはみ出しているかチェック
        # （正確な判定は難しいため、テキスト量で簡易判定）
        total_chars = sum(len(p.text) for p in text_frame.paragraphs)
        
        # シェイプのサイズに対してテキストが多すぎる場合
        shape_area = (shape.width / 914400) * (shape.height / 914400)  # 平方インチ
        chars_per_sqin = total_chars / shape_area if shape_area > 0 else 0
        
        # 経験的な閾値: 1平方インチあたり200文字以上は要注意
        if chars_per_sqin > 200:
            issues.append({
                'shape': shape,
                'text_length': total_chars,
                'area': shape_area,
                'density': chars_per_sqin,
                'message': f"テキストが多すぎる可能性: {total_chars}文字 / {shape_area:.2f}平方インチ"
            })
    
    return issues


def get_template_layouts_info(template_path):
    """
    テンプレートの各スライドをレイアウトカタログとして分析
    
    Args:
        template_path: テンプレートファイルのパス
    
    Returns:
        list: 各スライドの情報を含む辞書のリスト
    """
    prs = Presentation(template_path)
    layouts_info = []
    
    for i, slide in enumerate(prs.slides):
        info = {
            'slide_index': i,
            'layout_name': slide.slide_layout.name,
            'text_shapes_count': 0,
            'image_shapes_count': 0,
            'decorative_shapes_count': 0,
            'has_title': False,
            'has_content': False,
            'title_text': '',
            'content_preview': ''
        }
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                info['text_shapes_count'] += 1
                
                if is_title_shape(shape):
                    info['has_title'] = True
                    info['title_text'] = shape.text[:50]  # 最初の50文字
                elif is_content_shape(shape):
                    info['has_content'] = True
                    info['content_preview'] = shape.text[:100]  # 最初の100文字
            
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                info['image_shapes_count'] += 1
            else:
                info['decorative_shapes_count'] += 1
        
        # レイアウトタイプの推定
        if info['has_title'] and not info['has_content']:
            info['suggested_use'] = 'タイトルスライド、セクション区切り'
        elif info['has_title'] and info['has_content']:
            if info['text_shapes_count'] >= 3:
                info['suggested_use'] = '複数セクション、グリッドレイアウト'
            else:
                info['suggested_use'] = '箇条書き、標準コンテンツ'
        elif info['image_shapes_count'] > 0:
            info['suggested_use'] = '画像中心、図表スライド'
        else:
            info['suggested_use'] = 'その他'
        
        layouts_info.append(info)
    
    return layouts_info


def remove_unused_slides(prs, slides_to_keep):
    """
    使用しないスライドを削除
    
    Args:
        prs: Presentationオブジェクト
        slides_to_keep: 保持するスライドのインデックスのリスト
    
    Returns:
        int: 削除されたスライド数
    """
    # 削除するスライドのインデックスを特定
    all_indices = set(range(len(prs.slides)))
    to_remove = all_indices - set(slides_to_keep)
    
    # 後ろから削除（インデックスのずれを防ぐため）
    removed_count = 0
    for idx in sorted(to_remove, reverse=True):
        try:
            rId = prs.slides._sldIdLst[idx].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[idx]
            removed_count += 1
        except Exception as e:
            print(f"スライド{idx}の削除に失敗: {e}")
    
    return removed_count


def adjust_text_size_to_fit(shape, min_size_pt=12):
    """
    テキストがシェイプに収まるようにフォントサイズを調整
    
    Args:
        shape: 調整対象のシェイプ
        min_size_pt: 最小フォントサイズ（ポイント）
    
    Returns:
        bool: 調整が必要だった場合True
    """
    if not shape.has_text_frame:
        return False
    
    text_frame = shape.text_frame
    adjusted = False
    
    # 各段落のフォントサイズをチェック
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is not None:
                current_size_pt = run.font.size.pt
                
                # フォントサイズが大きすぎる場合は縮小
                if current_size_pt > 24:
                    new_size = max(current_size_pt * 0.9, min_size_pt)
                    run.font.size = Pt(new_size)
                    adjusted = True
    
    return adjusted


# サンプル使用例
if __name__ == '__main__':
    import sys
    
    if len(sys.argv) < 2:
        print("Usage: python helpers.py <template.pptx>")
        sys.exit(1)
    
    template_path = sys.argv[1]
    
    print(f"テンプレート分析: {template_path}\n")
    
    # レイアウト情報を取得
    layouts = get_template_layouts_info(template_path)
    
    print("=" * 80)
    print("レイアウトカタログ")
    print("=" * 80)
    
    for layout in layouts:
        print(f"\nスライド {layout['slide_index'] + 1}:")
        print(f"  レイアウト名: {layout['layout_name']}")
        print(f"  テキストシェイプ: {layout['text_shapes_count']}")
        print(f"  画像: {layout['image_shapes_count']}")
        print(f"  装飾要素: {layout['decorative_shapes_count']}")
        print(f"  タイトル: {'あり' if layout['has_title'] else 'なし'}")
        print(f"  コンテンツ: {'あり' if layout['has_content'] else 'なし'}")
        print(f"  推奨用途: {layout['suggested_use']}")
        
        if layout['title_text']:
            print(f"  タイトル例: {layout['title_text']}")
    
    print("\n" + "=" * 80)
    
    # レイアウト使用状況の分析
    prs = Presentation(template_path)
    usage, issues = analyze_layout_usage(prs)
    
    print("\nレイアウト使用状況:")
    for layout_name, count in usage.items():
        print(f"  {layout_name}: {count}回使用")
    
    if issues:
        print("\n⚠️  警告:")
        for issue in issues:
            print(f"  {issue['message']}")
    else:
        print("\n✓ レイアウトの多様性: 問題なし")
