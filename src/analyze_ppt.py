from pptx import Presentation

prs = Presentation("samples/Presentation.pptx")

# slideのテキスト要素を抽出
def extract_texts(prs):
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)

    print(text_runs)


# slide layoutとそのplaceholderを抽出
def extract_layout(prs):
    for i, layout in enumerate(prs.slide_layouts):
        print(f'--- Layout {i}: {layout.name} ---')
        for ph in layout.placeholders:
            print(f'  idx={ph.placeholder_format.idx}, type={ph.placeholder_format.type}, name="{ph.name}"')


# slideに含まれるプレースホルダーを抽出
def extract_placeholders(prs):
    for i, slide in enumerate(prs.slides):
        print(f'slide{i}')
        for shape in slide.placeholders:
            print(f'  idx={shape.placeholder_format.idx}, type={shape.placeholder_format.type}, name="{shape.name}"')

# 各スライドで使用されているlayoutを抽出
def extract_slide_layouts(prs):
    for i, slide in enumerate(prs.slides):
        layout = slide.slide_layout
        print(f'Slide {i}: layout="{layout.name}"')


def remove_slide(prs, slide_index):
    """プレゼンテーションからスライドを削除する。

    Args:
        prs: Presentationオブジェクト
        slide_index: 削除するスライドの0始まりインデックス
    """
    # 1. プレゼンテーションのスライドID一覧から該当要素を取得
    sldIdLst = prs._element.sldIdLst
    target_sldId = sldIdLst.sldId_lst[slide_index]

    # 2. スライドID一覧から削除
    sldIdLst.remove(target_sldId)

    # 3. リレーションシップを削除（パッケージからスライドパーツも除去される）
    prs.part.drop_rel(target_sldId.rId)


if __name__ == "__main__":
    #extract_texts(prs)
    extract_layout(prs)
    #extract_placeholders(prs)
    #extract_slide_layouts(prs)
